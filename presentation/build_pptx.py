"""Build the Data Management 2 final-project presentation (.pptx).

Clean SRH exposé style (Lato, orange srh wordmark, thin rules, flat callouts),
content = the real weather + air quality pipeline (Open-Meteo, Spark, BigQuery,
dbt, Airflow) with the verified run results.

Run with uv (the project has no pyproject, but --no-project is safe either way):
    uv run --no-project --with python-pptx python presentation/build_pptx.py

Output: presentation/DM2_Weather_Pipeline.pptx
Writing rule: no em dashes (commas, colons, parentheses). En dash only in ranges.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# ---- SRH exposé palette ----------------------------------------------------
ORANGE = RGBColor(0xE6, 0x44, 0x15)
INK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x55, 0x55, 0x55)
RULE = RGBColor(0xE3, 0xE3, 0xE3)
RULE2 = RGBColor(0xBB, 0xBB, 0xBB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FAINT = RGBColor(0xFA, 0xF6, 0xF4)
CODEBG = RGBColor(0xF4, 0xF4, 0xF4)

FONT = "Lato"
MONO = "Courier New"

HERE = Path(__file__).resolve().parent
LOGO = str(HERE / "srh_logo.jpg")
LOGO_RATIO = 457 / 591

ND = "–"  # en dash for ranges
GE = "≥"

SW = Inches(13.333)
SH = Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


# ---- helpers ---------------------------------------------------------------
def _no_shadow(shape):
    el = shape._element
    style = el.find(qn("p:style"))
    if style is not None:
        el.remove(style)
    spPr = el.spPr
    if spPr.find(qn("a:effectLst")) is None:
        spPr.append(parse_xml('<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'))


def slide():
    return prs.slides.add_slide(BLANK)


def R(txt, size=14, color=INK, bold=False, italic=False, font=FONT):
    return (txt, size, color, bold, italic, font)


def text(s, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6, line_spacing=1.12):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (txt, size, color, bold, italic, font) in para:
            run = p.add_run()
            run.text = txt
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.bold = bold
            run.font.italic = italic
            run.font.name = font
    return tb


def rule(s, x, y, w, color=RULE, pt=1.0):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Pt(pt))
    r.fill.solid()
    r.fill.fore_color.rgb = color
    r.line.fill.background()
    _no_shadow(r)
    return r


def box(s, x, y, w, h, fill=None, line=None, pt=1.0, rounded=False):
    shp = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    b = s.shapes.add_shape(shp, x, y, w, h)
    if fill is None:
        b.fill.background()
    else:
        b.fill.solid()
        b.fill.fore_color.rgb = fill
    if line is None:
        b.line.fill.background()
    else:
        b.line.color.rgb = line
        b.line.width = Pt(pt)
    _no_shadow(b)
    return b


def arrow(s, x, y, w, h, color=ORANGE):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    a.fill.solid()
    a.fill.fore_color.rgb = color
    a.line.fill.background()
    _no_shadow(a)
    return a


def logo_small(s):
    w = Inches(0.95)
    s.shapes.add_picture(LOGO, Inches(0.6), Inches(0.46), width=w, height=Emu(int(w * LOGO_RATIO)))


def pagenum(s, n):
    text(s, Inches(12.5), Inches(7.04), Inches(0.7), Inches(0.3), [[R(str(n), 10, RULE2)]],
         align=PP_ALIGN.RIGHT)


def header(s, num, title, n):
    logo_small(s)
    runs = []
    if num:
        runs.append(R(num + "   ", 27, ORANGE, bold=True))
    runs.append(R(title, 27, INK, bold=True))
    text(s, Inches(0.62), Inches(1.42), Inches(12.1), Inches(0.7), [runs])
    rule(s, Inches(0.64), Inches(2.18), Inches(12.05), color=RULE, pt=1.2)
    pagenum(s, n)
    return Inches(2.5)


def bullets(s, x, y, w, items, size=15, gap=9, lead=1.16):
    paras = []
    for txt, lvl in items:
        if lvl == 0:
            paras.append([R("•  ", size, ORANGE, bold=True), R(txt, size, INK)])
        else:
            paras.append([R("      –  ", size - 1, ORANGE), R(txt, size - 1, GRAY)])
    return text(s, x, y, w, Inches(4.4), paras, space_after=gap, line_spacing=lead)


def flow(s, y, items, box_h=1.15, fontsize=13, sub=11):
    """Horizontal flow of outlined boxes with orange arrows between them."""
    n = len(items)
    gap = Inches(0.2)
    left = Inches(0.7)
    total = int(SW) - 2 * int(left)
    bw = (total - int(gap) * (n - 1)) // n
    bh = Inches(box_h)
    for i, (t, st) in enumerate(items):
        bx = int(left) + (bw + int(gap)) * i
        box(s, Emu(bx), y, Emu(bw), bh, fill=WHITE, line=ORANGE, pt=1.5, rounded=True)
        text(s, Emu(bx + Emu(int(Inches(0.08)))), y + Inches(0.16), Emu(bw - int(Inches(0.16))), Inches(0.6),
             [[R(seg, fontsize, INK, bold=True)] for seg in t.split("\n")],
             align=PP_ALIGN.CENTER, space_after=0, line_spacing=1.0)
        text(s, Emu(bx + Emu(int(Inches(0.08)))), y + Inches(0.66), Emu(bw - int(Inches(0.16))), Inches(0.45),
             [[R(seg, sub, GRAY)] for seg in st.split("\n")],
             align=PP_ALIGN.CENTER, space_after=0, line_spacing=1.0)
        if i < n - 1:
            ax = bx + bw + (int(gap) - int(Inches(0.16))) // 2
            arrow(s, Emu(ax), y + Emu(int(bh)) // 2 - Inches(0.08), Inches(0.16), Inches(0.16))


def codebox(s, x, y, w, h, lines, caption=None):
    box(s, x, y, w, h, fill=CODEBG)
    top = y + Inches(0.18)
    if caption:
        text(s, x + Inches(0.25), top, w - Inches(0.5), Inches(0.3),
             [[R(caption, 10.5, GRAY, italic=True)]])
        top = top + Inches(0.34)
    text(s, x + Inches(0.25), top, w - Inches(0.5), h - Inches(0.4),
         [[R(ln, 11, INK, font=MONO)] for ln in lines], space_after=1, line_spacing=1.05)


# ===========================================================================
# Slide 1: Title
# ===========================================================================
s = slide()
lw = Inches(1.55)
s.shapes.add_picture(LOGO, Emu(int((int(SW) - int(lw)) / 2)), Inches(1.05), width=lw,
                     height=Emu(int(int(lw) * LOGO_RATIO)))
text(s, Inches(1.0), Inches(2.7), Inches(11.33), Inches(0.4),
     [[R("D A T A   M A N A G E M E N T   2      ·      F I N A L   P R O J E C T", 14, GRAY, bold=True)]],
     align=PP_ALIGN.CENTER)
text(s, Inches(1.0), Inches(3.25), Inches(11.33), Inches(0.9),
     [[R("Real-Time Weather and Air Quality Data Pipeline", 30, INK, bold=True)]],
     align=PP_ALIGN.CENTER, line_spacing=1.04)
text(s, Inches(1.0), Inches(4.15), Inches(11.33), Inches(0.5),
     [[R("Open-Meteo   ·   Spark   ·   BigQuery   ·   dbt   ·   Airflow", 17, ORANGE, bold=True)]],
     align=PP_ALIGN.CENTER)
rule(s, Emu(int((int(SW) - int(Inches(1.4))) / 2)), Inches(5.0), Inches(1.4), color=ORANGE, pt=3.0)
text(s, Inches(1.0), Inches(5.5), Inches(11.33), Inches(1.5),
     [[R("Devendra Singh Dhakad", 16, INK, bold=True),
       R("      ·      Matriculation No. 100004684", 16, GRAY)],
      [R("Data Management 2: Data Curation and Data Management", 14, GRAY)],
      [R("Examiner: Prof. Binh Vu", 14, GRAY)],
      [R("SRH University", 14, GRAY)]],
     align=PP_ALIGN.CENTER, space_after=7)
text(s, Inches(1.0), Inches(6.95), Inches(11.33), Inches(0.35),
     [[R("June 2026", 12, RULE2)]], align=PP_ALIGN.CENTER)


# ===========================================================================
# Slide 2: Project goal
# ===========================================================================
s = slide()
y = header(s, "1", "Project goal", 2)
bullets(s, Inches(0.7), y, Inches(11.9), [
    ("Build a complete, automated data pipeline that collects live weather and air "
     "quality for ten European cities, cleans and models the data, and serves "
     "analytics-ready tables in BigQuery.", 0),
    ("It demonstrates the full data-curation lifecycle: real-time extraction, "
     "distributed cleaning, warehouse loading, SQL transformation with tests, and "
     "hands-off scheduling.", 0),
], size=16, gap=12)
stats = [("2", "data sources\n(one real-time)"), ("10", "European cities\npolled live"),
         ("5 min", "automatic\nrun interval"), ("6 / 6", "exam requirements\nmet")]
sw_b = Inches(2.78)
gap = Inches(0.32)
x0 = Inches(0.7)
sy = Inches(4.45)
for i, (big, lab) in enumerate(stats):
    bx = int(x0) + (int(sw_b) + int(gap)) * i
    box(s, Emu(bx), sy, sw_b, Inches(1.85), fill=FAINT)
    box(s, Emu(bx), sy, sw_b, Inches(0.08), fill=ORANGE)
    text(s, Emu(bx), sy + Inches(0.32), sw_b, Inches(0.8),
         [[R(big, 36, ORANGE, bold=True)]], align=PP_ALIGN.CENTER)
    text(s, Emu(bx + int(Inches(0.1))), sy + Inches(1.18), Emu(int(sw_b) - int(Inches(0.2))), Inches(0.6),
         [[R(seg, 12, GRAY)] for seg in lab.split("\n")], align=PP_ALIGN.CENTER, space_after=0)


# ===========================================================================
# Slide 3: Architecture
# ===========================================================================
s = slide()
y = header(s, "2", "Architecture and data flow", 3)
flow(s, Inches(2.9), [
    ("Open-Meteo\nAPIs", "weather +\nair quality"),
    ("producer.py", "extract to\nJSON"),
    ("Spark", "clean to\nParquet"),
    ("load_bigquery", "load to\nraw table"),
    ("dbt", "model +\ntest"),
    ("Mart\ntables", "analytics\nready"),
], box_h=1.4)
box(s, Inches(0.7), Inches(4.75), Inches(11.9), Inches(0.6), fill=ORANGE)
text(s, Inches(0.7), Inches(4.75), Inches(11.9), Inches(0.6),
     [[R("Apache Airflow runs every step, every 5 minutes (cron fallback: run_pipeline.sh)",
         14, WHITE, bold=True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, Inches(0.7), Inches(5.7), Inches(11.9), Inches(0.9),
     [[R("Two isolated Python 3.11 environments keep dependencies apart: ", 13, INK, bold=True),
       R(".venv-pipeline (Spark, dbt, BigQuery) and airflow/.venv-airflow (Airflow 2.10). "
         "Airflow tasks shell into the pipeline environment.", 13, GRAY)]], line_spacing=1.15)


# ===========================================================================
# Slide 4: Data sources
# ===========================================================================
s = slide()
y = header(s, "3", "Data sources: two, one real-time", 4)
box(s, Inches(0.7), y, Inches(5.8), Inches(3.4), fill=FAINT)
box(s, Inches(0.7), y, Inches(0.1), Inches(3.4), fill=ORANGE)
text(s, Inches(1.0), y + Inches(0.22), Inches(5.3), Inches(3.0),
     [[R("Source 1: Open-Meteo  (real-time)", 15, ORANGE, bold=True)],
      [R("Two free REST APIs, no key. Weather (temperature, humidity, wind, code) and "
         "air quality (PM10, PM2.5, European AQI, ozone).", 13.5, INK)],
      [R("Polled for 10 cities every 5 minutes, so the landing zone is a live, growing "
         "stream of observations.", 13.5, INK)]], space_after=8, line_spacing=1.14)
box(s, Inches(6.85), y, Inches(5.75), Inches(3.4), fill=FAINT)
box(s, Inches(6.85), y, Inches(0.1), Inches(3.4), fill=ORANGE)
text(s, Inches(7.15), y + Inches(0.22), Inches(5.25), Inches(3.0),
     [[R("Source 2: City reference  (static)", 15, ORANGE, bold=True)],
      [R("A curated CSV of the 10 cities (country, country name, latitude, longitude, "
         "population, timezone), loaded as a dbt seed.", 13.5, INK)],
      [R("Reference data that rarely changes. It enriches every reading and anchors a "
         "referential-integrity test.", 13.5, INK)]], space_after=8, line_spacing=1.14)
text(s, Inches(0.7), y + Inches(3.65), Inches(11.9), Inches(0.5),
     [[R("The two sources join on city: a fast-moving stream combined with slow-moving "
         "reference data, a classic curation pattern.", 13.5, INK, italic=True)]])


# ===========================================================================
# Slide 5: Extract
# ===========================================================================
s = slide()
y = header(s, "4", "Extract: producer.py", 5)
bullets(s, Inches(0.7), y, Inches(5.7), [
    ("Calls both Open-Meteo endpoints for each city and writes one newline-delimited "
     "JSON file per run into the landing zone.", 0),
    ("timezone=UTC keeps every observation timestamp comparable.", 0),
    ("Per-city try and except: one failed city does not abort the whole run.", 0),
    ("The nesting is deliberate, so Spark does the real flattening downstream.", 0),
], size=14, gap=10)
codebox(s, Inches(6.6), y, Inches(6.0), Inches(4.0), [
    "{",
    '  "city": "Berlin", "country": "DE",',
    '  "weather": {',
    '     "time": "2026-06-14T15:45",',
    '     "temperature_2m": 15.9,',
    '     "relative_humidity_2m": 64,',
    '     "wind_speed_10m": 20.8 },',
    '  "air_quality": {',
    '     "pm2_5": 2.4, "european_aqi": 26 },',
    '  "ingested_at": "2026-06-14T15:50:19Z"',
    "}",
], caption="data/landing/readings_*.json  (one line per city)")


# ===========================================================================
# Slide 6: Spark cleaning
# ===========================================================================
s = slide()
y = header(s, "5", "Clean with Spark", 6)
steps = [
    "Flatten the nested weather and air-quality blocks",
    "Cast every field to its proper type",
    "Parse the time string into a real timestamp",
    "Drop rows with null city, time or temperature",
    "Deduplicate by (city, observed_at) in the batch",
    "Derive aqi_category from the European AQI bands",
]
for i, st in enumerate(steps):
    yy = y + Emu(int(Inches(0.6)) * i)
    box(s, Inches(0.7), yy, Inches(0.42), Inches(0.42), fill=ORANGE, rounded=True)
    text(s, Inches(0.7), yy, Inches(0.42), Inches(0.42), [[R(str(i + 1), 14, WHITE, bold=True)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(1.28), yy, Inches(5.1), Inches(0.45), [[R(st, 13, INK)]], anchor=MSO_ANCHOR.MIDDLE)
codebox(s, Inches(6.6), y, Inches(6.0), Inches(3.7), [
    "cleaned = (batch_df.select(",
    "    F.col('city'),",
    "    F.to_timestamp(F.col('weather.time'),",
    "        \"yyyy-MM-dd'T'HH:mm\")",
    "      .alias('observed_at'),",
    "    F.col('weather.temperature_2m'))",
    "  .where(F.col('temperature_2m').isNotNull())",
    "  .dropDuplicates(['city','observed_at']))",
], caption="Spark Structured Streaming, Trigger.AvailableNow")


# ===========================================================================
# Slide 7: BigQuery load
# ===========================================================================
s = slide()
y = header(s, "6", "Load into BigQuery", 7)
bullets(s, Inches(0.7), y, Inches(6.0), [
    ("load_bigquery.py reads all cleaned Parquet and loads it into "
     "raw.weather_readings with WRITE_TRUNCATE.", 0),
    ("Idempotent: running it twice gives the same table, no duplicates, so Airflow "
     "can safely retry.", 0),
    ("Uses the google-cloud-bigquery client only. No gcloud CLI; auth via a service "
     "account key.", 0),
    ("Verified: 30 rows loaded into the raw table.", 0),
], size=14, gap=10)
box(s, Inches(6.9), y, Inches(5.7), Inches(3.9), fill=FAINT)
box(s, Inches(6.9), y, Inches(0.1), Inches(3.9), fill=ORANGE)
text(s, Inches(7.2), y + Inches(0.22), Inches(5.2), Inches(3.5),
     [[R("raw.weather_readings", 14, ORANGE, bold=True, font=MONO)],
      [R("city, country            STRING", 12.5, INK, font=MONO)],
      [R("latitude, longitude      FLOAT", 12.5, INK, font=MONO)],
      [R("observed_at              TIMESTAMP", 12.5, INK, font=MONO)],
      [R("temperature_2m           FLOAT", 12.5, INK, font=MONO)],
      [R("relative_humidity_2m     INTEGER", 12.5, INK, font=MONO)],
      [R("wind_speed_10m           FLOAT", 12.5, INK, font=MONO)],
      [R("pm10, pm2_5, ozone       FLOAT", 12.5, INK, font=MONO)],
      [R("european_aqi             INTEGER", 12.5, INK, font=MONO)],
      [R("aqi_category             STRING", 12.5, INK, font=MONO)],
      [R("ingested_at              TIMESTAMP", 12.5, INK, font=MONO)]],
     space_after=4, line_spacing=1.1)


# ===========================================================================
# Slide 8: dbt transformation
# ===========================================================================
s = slide()
y = header(s, "7", "Transform with dbt (ELT)", 8)
flow(s, Inches(2.55), [
    ("raw.weather_readings", "source"),
    ("stg_weather", "staging (view)"),
    ("fct_weather_readings", "mart (table)"),
    ("agg_* KPIs", "mart (table)"),
], box_h=1.0, fontsize=12, sub=10.5)
text(s, Inches(0.7), Inches(3.75), Inches(11.9), Inches(0.5),
     [[R("The static seed (stg_cities) joins into fct_weather_readings to add country "
         "name and population.", 12.5, GRAY, italic=True)]])
box(s, Inches(0.7), Inches(4.35), Inches(5.8), Inches(2.45), fill=FAINT)
box(s, Inches(0.7), Inches(4.35), Inches(0.1), Inches(2.45), fill=ORANGE)
text(s, Inches(1.0), Inches(4.55), Inches(5.3), Inches(2.1),
     [[R("Why ELT, not ETL", 14, ORANGE, bold=True)],
      [R("Spark does only light cleaning. All business logic then runs as dbt SQL inside "
         "BigQuery: warehouse scale, version-controlled SQL, lineage, and built-in tests.",
         13, INK)]], space_after=6, line_spacing=1.14)
box(s, Inches(6.85), Inches(4.35), Inches(5.75), Inches(2.45), fill=FAINT)
box(s, Inches(6.85), Inches(4.35), Inches(0.1), Inches(2.45), fill=ORANGE)
text(s, Inches(7.15), Inches(4.55), Inches(5.25), Inches(2.1),
     [[R("The models", 14, ORANGE, bold=True)],
      [R("staging:  stg_weather, stg_cities", 12.5, INK, font=MONO)],
      [R("marts:    fct_weather_readings", 12.5, INK, font=MONO)],
      [R("          agg_weather_by_city", 12.5, INK, font=MONO)],
      [R("          agg_air_quality_by_country", 12.5, INK, font=MONO)],
      [R("5 models built successfully.", 12.5, INK)]], space_after=5, line_spacing=1.14)


# ===========================================================================
# Slide 9: Data quality
# ===========================================================================
s = slide()
y = header(s, "8", "Data quality: 17 tests, all pass", 9)
tests = [
    ("not_null", "Key fields are never empty: city, observed_at, reading_id."),
    ("unique", "No duplicate keys: reading_id and city are unique in the marts."),
    ("accepted_values", "aqi_category is always one of the six European AQI bands."),
    ("relationships", "Every reading's city exists in the static city reference."),
]
for i, (name, desc) in enumerate(tests):
    yy = y + Emu(int(Inches(0.92)) * i)
    box(s, Inches(0.7), yy, Inches(2.7), Inches(0.72), fill=FAINT)
    box(s, Inches(0.7), yy, Inches(0.08), Inches(0.72), fill=ORANGE)
    text(s, Inches(0.7), yy, Inches(2.7), Inches(0.72), [[R(name, 15, ORANGE, bold=True, font=MONO)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(3.65), yy, Inches(8.9), Inches(0.72), [[R(desc, 14, INK)]], anchor=MSO_ANCHOR.MIDDLE)
text(s, Inches(0.7), Inches(6.5), Inches(11.9), Inches(0.5),
     [[R("dbt test fails the Airflow task if any check fails, so bad data is caught before "
         "anyone trusts the KPIs.", 13.5, GRAY, italic=True)]])


# ===========================================================================
# Slide 10: Automation (Airflow) with real evidence
# ===========================================================================
s = slide()
y = header(s, "9", "Automation: Apache Airflow", 10)
flow(s, Inches(2.55), [
    ("ingest", ""), ("spark_clean", ""), ("load_bigquery", ""),
    ("dbt_seed", ""), ("dbt_run", ""), ("dbt_test", ""),
], box_h=0.85, fontsize=12.5)
text(s, Inches(0.7), Inches(3.6), Inches(11.9), Inches(0.4),
     [[R("schedule = '*/5 * * * *'      catchup = False      max_active_runs = 1      retries = 1",
         12.5, GRAY, font=MONO)]], align=PP_ALIGN.CENTER)
box(s, Inches(0.7), Inches(4.25), Inches(11.9), Inches(1.95), fill=FAINT)
box(s, Inches(0.7), Inches(4.25), Inches(0.1), Inches(1.95), fill=ORANGE)
text(s, Inches(1.0), Inches(4.45), Inches(11.4), Inches(1.6),
     [[R("Verified run:  ", 14, ORANGE, bold=True),
       R("the DAG executed end to end through Airflow. DagRun state = success; all six "
         "tasks (ingest, spark_clean, load_bigquery, dbt_seed, dbt_run, dbt_test) returned "
         "success.", 14, INK)],
      [R("dbt_seed loaded the cities, dbt_run built 5 models, dbt_test passed all 17 "
         "checks against BigQuery.", 13, GRAY)]], space_after=8, line_spacing=1.16)
text(s, Inches(0.7), Inches(6.45), Inches(11.9), Inches(0.4),
     [[R("run_pipeline.sh runs the same six-step chain and can be scheduled with cron as a "
         "fallback.", 12.5, GRAY, italic=True)]])


# ===========================================================================
# Slide 11: Sample results
# ===========================================================================
s = slide()
y = header(s, "10", "Sample results from the marts", 11)
text(s, Inches(0.7), y - Inches(0.05), Inches(7.0), Inches(0.4),
     [[R("Air quality by country  (average European AQI, lower is better)", 13.5, GRAY, italic=True)]])
kpi = [
    ("Country", "Avg PM2.5", "Avg EAQI"),
    ("Italy", "4.7", "45.0"),
    ("Portugal", "9.7", "41.0"),
    ("Spain", "6.1", "40.0"),
    ("France", "4.4", "34.0"),
    ("Germany", "3.0", "29.5"),
    ("Poland", "4.3", "26.5"),
]
ty = y + Inches(0.45)
for i, row in enumerate(kpi):
    head = i == 0
    if head:
        box(s, Inches(0.7), ty, Inches(6.0), Inches(0.5), fill=ORANGE)
    elif i % 2 == 0:
        box(s, Inches(0.7), ty, Inches(6.0), Inches(0.46), fill=FAINT)
    cols = [Inches(0.9), Inches(3.4), Inches(5.0)]
    for c, val in enumerate(row):
        col = WHITE if head else INK
        text(s, cols[c], ty + Inches(0.06), Inches(2.4), Inches(0.4),
             [[R(val, 13, col, bold=head)]])
    ty = ty + (Inches(0.5) if head else Inches(0.46))
box(s, Inches(7.1), y + Inches(0.4), Inches(5.5), Inches(3.4), fill=FAINT)
box(s, Inches(7.1), y + Inches(0.4), Inches(0.1), Inches(3.4), fill=ORANGE)
text(s, Inches(7.4), y + Inches(0.6), Inches(5.0), Inches(3.0),
     [[R("Highlights", 15, ORANGE, bold=True)],
      [R("Hottest city:  Madrid, 31.8 C average", 13.5, INK)],
      [R("Coolest city:  Warsaw, 13.25 C average", 13.5, INK)],
      [R("Worst air quality:  Italy (EAQI 45)", 13.5, INK)],
      [R("Best air quality:  Poland (EAQI 26.5)", 13.5, INK)],
      [R("Counts grow every 5 minutes as new hourly observations arrive, so the averages "
         "become richer over time.", 13, GRAY)]], space_after=9, line_spacing=1.14)


# ===========================================================================
# Slide 12: Requirement mapping
# ===========================================================================
s = slide()
y = header(s, "11", "How each requirement is met", 12)
mapping = [
    ("Requirement", "Where it is implemented"),
    ("1. Two sources, one real-time", "Open-Meteo APIs (real-time) + city reference seed (static)"),
    ("2. Extract, clean, load to BigQuery", "producer.py, spark_clean.py, load_bigquery.py"),
    ("3. Transformation with dbt", "seed + staging + marts, all SQL inside BigQuery (ELT)"),
    ("4. Include Spark", "Spark Structured Streaming cleaning job"),
    ("5. Runs automatically", "Airflow DAG every 5 min (cron fallback: run_pipeline.sh)"),
    ("6. Presentation", "this deck: sources, cleaning, transformation explained"),
]
ty = y
for i, (a, b) in enumerate(mapping):
    head = i == 0
    rh = Inches(0.6) if head else Inches(0.62)
    if head:
        box(s, Inches(0.7), ty, Inches(11.9), rh, fill=ORANGE)
    elif i % 2 == 0:
        box(s, Inches(0.7), ty, Inches(11.9), rh, fill=FAINT)
    col = WHITE if head else INK
    text(s, Inches(0.9), ty + Inches(0.12), Inches(4.3), Inches(0.45), [[R(a, 13.5, col, bold=True)]])
    text(s, Inches(5.3), ty + Inches(0.12), Inches(7.1), Inches(0.45), [[R(b, 13.5, col, bold=head)]])
    ty = ty + rh


out = HERE / "DM2_Weather_Pipeline.pptx"
prs.save(str(out))
print(f"Saved {out} with {len(prs.slides._sldIdLst)} slides")
